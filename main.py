from PyQt5.QtWidgets import QApplication, QLabel, QPushButton, QVBoxLayout, \
    QWidget, QDesktopWidget, QSizePolicy, QFileDialog, QMessageBox
from PyQt5.QtGui import QDropEvent, QIcon, QDragEnterEvent
from PyQt5.QtCore import Qt, pyqtSignal, QThread
from io import BytesIO
from pptx import Presentation
from PIL import Image
from gtts import gTTS
import numpy as np
import pytesseract
import os
import cv2
class DragAndDropLabel(QLabel):
    file_dropped = pyqtSignal(str)

    def __init__(self, parent=None):
        super(DragAndDropLabel, self).__init__(parent)
        self.setAlignment(Qt.AlignCenter)
        self.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        self.setText("Arrastra tu presentación aquí")
        self.setStyleSheet("border: 5px dashed #000814; color: #032B43; font-size: 35px;"
                           " font-weight: bold; padding: 40px; margin-top: 20px; background-color: #3F88C5; ")

        self.setAcceptDrops(True)

    def dragEnterEvent(self, event: QDragEnterEvent):
        mime_data = event.mimeData()
        if mime_data.hasUrls() and len(mime_data.urls()) == 1:
            event.acceptProposedAction()

    def dropEvent(self, event: QDropEvent):
        file_path = event.mimeData().urls()[0].toLocalFile()
        self.file_dropped.emit(file_path)

class InitialScreen(QWidget):
    file_dropped = pyqtSignal(str)

    def __init__(self):
        super().__init__()
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()
        pantalla = QDesktopWidget().availableGeometry()
        self.setWindowTitle("Conversión Texto-Audio")
        self.setGeometry(100, 100, 1280, 700)
        self.setWindowIcon(QIcon('esime_original.ico'))
        self.setStyleSheet("background-color: #669bbc;")

        self.titulo1 = QLabel('Convierte un PowerPoint a MP3', self)
        self.titulo1.setStyleSheet("color: #000814; font-size: 40px; font-weight: bold;")
        layout.addWidget(self.titulo1, alignment=Qt.AlignCenter)

        self.drag_and_drop_label = DragAndDropLabel(self)
        self.drag_and_drop_label.file_dropped.connect(self.file_dropped.emit)
        layout.addWidget(self.drag_and_drop_label, alignment=Qt.AlignTop)
        height_percentage = 0.6
        label_height = int(pantalla.height() * height_percentage)
        self.drag_and_drop_label.setMinimumHeight(label_height)
        self.drag_and_drop_label.setMaximumHeight(label_height)

        self.setLayout(layout)

    def handle_file_dropped(self, file_path):
        if file_path.lower().endswith(('.pptx', '.ppt')):
            self.file_dropped.emit(file_path)
        else:
            QMessageBox.warning(self, "Advertencia", "Formato de archivo no válido. Por favor, selecciona un archivo PowerPoint.")

class AudioConversionThread(QThread):
    finished = pyqtSignal(str)

    def __init__(self, pptx_path, base_audio_filename):
        super().__init__()
        self.pptx_path = pptx_path
        self.base_audio_filename = base_audio_filename

    def run(self):
        try:
            pptx_info = self.extract_pptx_info_with_ocr(self.pptx_path)
            self.perform_ocr_on_images(pptx_info)

            base_audio_filename = "presentacion"
            self.generate_audio_file(pptx_info, base_audio_filename)
            self.finished.emit("success")
        except Exception as e:
            self.finished.emit(str(e))

    def extract_table_text(self, table):
        table_text = ""
        for row in table.rows:
            for cell in row.cells:
                table_text += cell.text + " "
            table_text += "\n"
        return table_text

    def extract_pptx_info_with_ocr(self, pptx_path):
        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        prs = Presentation(pptx_path)
        pptx_info = {'total_slides': len(prs.slides), 'slides': []}

        for slide in prs.slides:
            slide_info = {'title': "", 'text': "", 'images': [], 'tables': []}

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        if not slide_info['title']:
                            slide_info['title'] = text
                        else:
                            slide_info['text'] += text + "\n"
                elif shape.shape_type == 13:  # Shape type for images
                    img_stream = shape.image.blob
                    img = Image.open(BytesIO(img_stream))
                    slide_info['images'].append({'image': img, 'ocr_text': ''})
                elif shape.has_table:
                    table_text = self.extract_table_text(shape.table)
                    slide_info['tables'].append(table_text)

            pptx_info['slides'].append(slide_info)

        return pptx_info

    def perform_ocr_on_images(self, pptx_info):
        for slide in pptx_info['slides']:
            for img_data in slide['images']:
                img_array = cv2.cvtColor(np.array(img_data['image']), cv2.COLOR_RGB2BGR)
                text = pytesseract.image_to_string(img_array)
                img_data['ocr_text'] = text

    def save_text_to_audio(self, text, base_filename):
        counter = 0
        audio_filename = f"{base_filename}.mp3"

        while os.path.exists(audio_filename):
            counter += 1
            audio_filename = f"{base_filename}_{counter}.mp3"

        tts = gTTS(text=text, lang='es')
        tts.save(audio_filename)

    def generate_audio_file(self, pptx_info, base_filename):
        audio_text = ""

        for i, slide in enumerate(pptx_info['slides']):
            audio_text += f"Diapositiva {i + 1} - Título: {slide['title']}"
            if slide['text']:
                audio_text += f"Contenido: {slide['text']}"
            for j, table_text in enumerate(slide['tables']):
                audio_text += f"Tabla {j + 1} de la diapositiva {i + 1} - Contenido: {table_text}"
            for j, img_data in enumerate(slide['images']):
                if img_data['ocr_text']:
                    audio_text += f"Imagen {j + 1} de la diapositiva {i + 1} - Texto: {img_data['ocr_text']}"
                else:
                    audio_text += f"Imagen {j + 1} de la diapositiva {i + 1} - Sin contenido"

        self.save_text_to_audio(audio_text, base_filename)

class PowerPointToAudioConverter(QWidget):
    def __init__(self):
        super().__init__()
        self.pptx_path = None
        self.audio_conversion_thread = None
        self.init_ui()

    def init_ui(self):
        layout = QVBoxLayout()

        pantalla = QDesktopWidget().availableGeometry()
        self.setWindowTitle("Conversión Texto-Audio - Conversión")
        self.setGeometry(100, 100, 800, 600)
        self.setWindowIcon(QIcon('esime_original.ico'))
        self.setStyleSheet("background-color: #4059AD;")

        self.titulo2 = QLabel('Detalles de tu archivo', self)
        self.titulo2.setStyleSheet("color: #EFF2F1; font-size: 40px; font-weight: bold;")
        layout.addWidget(self.titulo2, alignment=Qt.AlignCenter)

        self.boton_convertir = QPushButton("Convertir a Audio", self)
        self.boton_convertir.setStyleSheet("background-color: #EFF2F1; color: black; font-size: 20px; font-weight: bold; border: 2px solid white;")
        self.boton_convertir.clicked.connect(self.convertir_a_audio)
        self.boton_convertir.setFixedSize(200, 50)
        self.boton_convertir.setMaximumWidth(200)
        self.boton_convertir.setContentsMargins(10, 10, 10, 10)
        layout.addWidget(self.boton_convertir, alignment=Qt.AlignCenter)

        self.etiqueta_seleccionado = QLabel("", self)
        self.etiqueta_seleccionado.setStyleSheet("color: white; font-size: 15px; font-weight: bold;")
        layout.addWidget(self.etiqueta_seleccionado, alignment=Qt.AlignTop)

        self.nombre_archivo_label = QLabel("", self)
        self.nombre_archivo_label.setStyleSheet("color: #2980B9; font-size: 15px; font-weight: bold;")
        layout.addWidget(self.nombre_archivo_label, alignment=Qt.AlignTop)

        self.vista_previa_label = QLabel(self)
        layout.addWidget(self.vista_previa_label, alignment=Qt.AlignTop)

        self.setLayout(layout)

    def obtener_numero_diapositivas(self, file_path):
        try:
            presentation = Presentation(file_path)
            numero_diapositivas = len(presentation.slides)
            return numero_diapositivas
        except Exception as e:
            mensaje = f"Error al obtener el número de diapositivas"
            return None

    def show_conversion_screen(self, file_path):
        self.etiqueta_seleccionado.setText("Archivo seleccionado:")
        self.nombre_archivo_label.setText(os.path.basename(file_path))
        self.nombre_archivo_label.setStyleSheet("color: #EFF2F1; font-size: 30px; font-weight: bold;")
        self.pptx_path = file_path

        numero_diapositivas = self.obtener_numero_diapositivas(file_path)
        mensaje = f"Número de diapositivas: {numero_diapositivas}"
        self.etiqueta_seleccionado.setText(mensaje)

        self.show()

    def convertir_a_audio(self):
        if not self.pptx_path:
            QMessageBox.warning(self, "Advertencia", "Primero selecciona un archivo PowerPoint.")
            return

        if self.audio_conversion_thread is None or not self.audio_conversion_thread.isRunning():
            self.audio_conversion_thread = AudioConversionThread(self.pptx_path, "presentacion")
            self.audio_conversion_thread.finished.connect(self.conversion_completed)
            self.audio_conversion_thread.start()
        else:
            QMessageBox.warning(self, "Advertencia", "La conversión de audio ya está en progreso.")

    def conversion_completed(self, result):
        if result == "success":
            QMessageBox.information(self, "Éxito", f"La conversión a audio ha terminado. Se ha guardado en el archivo presentacion.mp3")
        else:
            QMessageBox.warning(self, "Error", f"Error durante la conversión: {result}")

if __name__ == "__main__":
    app = QApplication([])
    main_window = InitialScreen()
    conversion_screen = PowerPointToAudioConverter()

    main_window.file_dropped.connect(conversion_screen.show_conversion_screen)

    main_window.show()
    app.exec_()
